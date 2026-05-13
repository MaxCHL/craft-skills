"""
pptx-craft primitive helpers.

Built on python-pptx. Each function takes a Slide object and lays out one
primitive on it. Recipes compose multiple primitives across slides.

Slide size is assumed 13.333" x 7.5" (16:9 widescreen).
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ----- Style tokens --------------------------------------------------------

@dataclass(frozen=True)
class Palette:
    bg: str
    ink: str
    accent: str
    muted: str = "888888"


@dataclass(frozen=True)
class Typography:
    headline_font: str
    body_font: str
    headline_size: int
    body_size: int


STYLES: dict[str, dict] = {
    "swiss": {
        "palette": Palette(bg="FFFFFF", ink="111111", accent="E63946", muted="888888"),
        "type": Typography("Inter", "Inter", 40, 14),
    },
    "editorial": {
        "palette": Palette(bg="F5F1EA", ink="1A1A1A", accent="8B2331", muted="6B6258"),
        "type": Typography("Source Serif Pro", "Source Serif Pro", 56, 13),
    },
    "bento": {
        "palette": Palette(bg="F2F2F2", ink="1D1D1F", accent="0066FF", muted="86868B"),
        "type": Typography("Inter", "Inter", 36, 13),
    },
    "brutalist": {
        "palette": Palette(bg="000000", ink="FFFFFF", accent="FFFF00", muted="666666"),
        "type": Typography("Impact", "Arial", 120, 14),
    },
    "big-type": {
        "palette": Palette(bg="0A1828", ink="F8F9FA", accent="FFA62B", muted="6C7A89"),
        "type": Typography("Inter", "Inter", 110, 14),
    },
    "data-dense": {
        "palette": Palette(bg="F8F5EE", ink="1A1A1A", accent="990000", muted="637478"),
        "type": Typography("Inter", "Inter", 24, 11),
    },
    "monograph": {
        "palette": Palette(bg="FAFAF7", ink="222222", accent="8B0000", muted="5A5A5A"),
        "type": Typography("EB Garamond", "EB Garamond", 28, 12),
    },
}


# ----- Low-level helpers ---------------------------------------------------

def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


def set_background(slide, hex_color: str) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(hex_color)
    bg.line.fill.background()
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(
    slide,
    text: str,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    *,
    font: str,
    size: int,
    color: str,
    bold: bool = False,
    italic: bool = False,
    align: str = "left",
    line_spacing: float = 1.15,
    tracking: float = 0.0,
):
    """Place a text box at absolute coordinates with style tokens."""
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    para = tf.paragraphs[0]
    para.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    para.line_spacing = line_spacing

    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    return tb


def add_rect(
    slide,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    *,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 0.75,
    rounded: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(
        shape_type, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(line_width)
    if rounded and hasattr(s, "adjustments"):
        # Soften the corner radius
        s.adjustments[0] = 0.08
    return s


def add_line(
    slide,
    x1_in: float,
    y1_in: float,
    x2_in: float,
    y2_in: float,
    *,
    color: str,
    width_pt: float = 1.0,
):
    from pptx.util import Emu
    line = slide.shapes.add_connector(
        1, Inches(x1_in), Inches(y1_in), Inches(x2_in), Inches(y2_in)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width_pt)
    return line


# ----- Primitives ---------------------------------------------------------

def primitive_manifesto(slide, text: str, style: str = "big-type", accent_word: str | None = None):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    # Vertical center slightly above optical midpoint (40% from top)
    add_text(
        slide, text,
        x_in=1.0, y_in=2.6, w_in=11.33, h_in=2.5,
        font=typ.headline_font, size=typ.headline_size,
        color=pal.ink,
        align="left" if style != "big-type" else "left",
        line_spacing=1.05,
    )


def primitive_pull_quote(slide, quote: str, attribution: str | None, style: str = "editorial"):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    if style == "swiss":
        # Left rule + no quote marks
        add_rect(slide, 1.0, 2.5, 0.04, 2.5, fill=pal.accent)
        add_text(
            slide, quote,
            x_in=1.3, y_in=2.5, w_in=10.5, h_in=2.5,
            font=typ.headline_font, size=28, color=pal.ink, line_spacing=1.3,
        )
        if attribution:
            add_text(
                slide, attribution.upper(),
                x_in=1.3, y_in=5.2, w_in=10.5, h_in=0.4,
                font=typ.body_font, size=10, color=pal.muted, tracking=200,
            )
    elif style == "brutalist":
        add_text(
            slide, quote,
            x_in=0.6, y_in=2.4, w_in=12.0, h_in=3.0,
            font=typ.headline_font, size=48, color=pal.ink, line_spacing=1.0,
        )
        if attribution:
            add_text(
                slide, "/ " + attribution,
                x_in=0.6, y_in=5.8, w_in=12.0, h_in=0.5,
                font="Courier New", size=12, color=pal.accent,
            )
    else:  # editorial
        add_text(
            slide, "“" + quote + "”",
            x_in=1.5, y_in=2.3, w_in=10.3, h_in=3.0,
            font=typ.headline_font, size=32, color=pal.ink,
            italic=True, line_spacing=1.3,
        )
        if attribution:
            add_text(
                slide, "— " + attribution,
                x_in=1.5, y_in=5.5, w_in=10.3, h_in=0.5,
                font=typ.body_font, size=12, color=pal.muted, italic=True,
            )


def primitive_big_stat(slide, number: str, label: str, context: str | None = None, style: str = "swiss"):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    # Number flush-left, oversized
    size = 200 if style == "big-type" else 140
    add_text(
        slide, number,
        x_in=0.8, y_in=1.5, w_in=8.0, h_in=3.5,
        font=typ.headline_font, size=size, color=pal.accent, line_spacing=1.0,
    )
    add_text(
        slide, label,
        x_in=0.8, y_in=5.2, w_in=10.0, h_in=0.6,
        font=typ.body_font, size=18, color=pal.ink,
    )
    if context:
        add_text(
            slide, context,
            x_in=0.8, y_in=5.9, w_in=10.0, h_in=1.0,
            font=typ.body_font, size=11, color=pal.muted, italic=True,
        )


def primitive_swiss_numbered_list(slide, items: Sequence[tuple[str, str]], title: str | None = None):
    tokens = STYLES["swiss"]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    y = 0.7
    if title:
        add_text(
            slide, title,
            x_in=0.8, y_in=y, w_in=12.0, h_in=0.7,
            font=typ.headline_font, size=28, color=pal.ink,
        )
        y += 1.0

    row_h = (6.5 - y) / max(len(items), 1)
    for num, text in items:
        # Hanging numeral
        add_text(
            slide, num,
            x_in=0.8, y_in=y, w_in=1.0, h_in=row_h,
            font=typ.headline_font, size=14, color=pal.accent,
        )
        add_text(
            slide, text,
            x_in=2.0, y_in=y, w_in=10.5, h_in=row_h,
            font=typ.body_font, size=14, color=pal.ink, line_spacing=1.4,
        )
        y += row_h


def primitive_bento_3x2(slide, cells: Sequence[dict], style: str = "bento"):
    """Each cell dict: {type, title, body, accent?: bool}.
    type ∈ {stat, quote, icon, code, photo, chart} — recipes just render text content;
    real images/charts get inserted via slide.shapes.add_picture in recipes."""
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    margin = 0.4
    gutter = 0.15
    grid_w = 13.333 - 2 * margin
    grid_h = 7.5 - 2 * margin
    cell_w = (grid_w - 2 * gutter) / 3
    cell_h = (grid_h - gutter) / 2

    for idx, cell in enumerate(cells[:6]):
        row, col = divmod(idx, 3)
        x = margin + col * (cell_w + gutter)
        y = margin + row * (cell_h + gutter)
        bg_color = pal.accent if cell.get("accent") else "FFFFFF"
        ink = "FFFFFF" if cell.get("accent") else pal.ink
        add_rect(slide, x, y, cell_w, cell_h, fill=bg_color, rounded=True)

        # Title (small)
        add_text(
            slide, cell.get("title", "").upper(),
            x_in=x + 0.25, y_in=y + 0.25, w_in=cell_w - 0.5, h_in=0.4,
            font=typ.body_font, size=10, color=ink, bold=True,
        )
        # Body (depends on type)
        body = cell.get("body", "")
        ctype = cell.get("type", "text")
        if ctype == "stat":
            add_text(
                slide, body,
                x_in=x + 0.25, y_in=y + 0.7, w_in=cell_w - 0.5, h_in=cell_h - 1.0,
                font=typ.headline_font, size=44, color=ink, line_spacing=1.0,
            )
        elif ctype == "quote":
            add_text(
                slide, "“" + body + "”",
                x_in=x + 0.25, y_in=y + 0.7, w_in=cell_w - 0.5, h_in=cell_h - 1.0,
                font=typ.body_font, size=13, color=ink, italic=True, line_spacing=1.4,
            )
        elif ctype == "code":
            add_text(
                slide, body,
                x_in=x + 0.25, y_in=y + 0.7, w_in=cell_w - 0.5, h_in=cell_h - 1.0,
                font="Consolas", size=10, color=ink, line_spacing=1.3,
            )
        else:
            add_text(
                slide, body,
                x_in=x + 0.25, y_in=y + 0.7, w_in=cell_w - 0.5, h_in=cell_h - 1.0,
                font=typ.body_font, size=13, color=ink, line_spacing=1.4,
            )


def primitive_data_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    accent_cell: tuple[int, int] | None = None,
    title: str | None = None,
    source: str | None = None,
    style: str = "data-dense",
):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    if title:
        add_text(
            slide, title,
            x_in=0.8, y_in=0.6, w_in=12.0, h_in=0.6,
            font=typ.headline_font, size=22, color=pal.ink,
        )

    from pptx.util import Inches as I
    cols = len(headers)
    rows_n = len(rows) + 1
    table_x, table_y, table_w, table_h = 0.8, 1.5, 11.7, 5.0
    table = slide.shapes.add_table(rows_n, cols, I(table_x), I(table_y), I(table_w), I(table_h)).table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0]
        run.font.name = typ.body_font
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = rgb(pal.ink)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(pal.bg)

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0]
            run.font.name = typ.body_font
            run.font.size = Pt(11)
            is_num = j > 0
            run.font.color.rgb = rgb(pal.ink)
            cell.fill.solid()
            if accent_cell and accent_cell == (i, j):
                cell.fill.fore_color.rgb = rgb(pal.accent)
                run.font.color.rgb = rgb(pal.bg)
            else:
                cell.fill.fore_color.rgb = rgb(pal.bg)
            para.alignment = PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT

    if source:
        add_text(
            slide, source,
            x_in=0.8, y_in=6.9, w_in=11.7, h_in=0.4,
            font=typ.body_font, size=9, color=pal.muted, italic=True,
        )


def primitive_chapter_opener(slide, kicker: str, title: str, meta: str | None = None, style: str = "editorial"):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    add_text(
        slide, kicker.upper(),
        x_in=1.0, y_in=2.5, w_in=11.0, h_in=0.4,
        font="Inter", size=11, color=pal.accent, tracking=200,
    )
    add_text(
        slide, title,
        x_in=1.0, y_in=3.0, w_in=11.0, h_in=3.0,
        font=typ.headline_font, size=64, color=pal.ink, line_spacing=1.05,
    )
    if meta:
        add_text(
            slide, meta,
            x_in=1.0, y_in=6.7, w_in=11.0, h_in=0.5,
            font="Inter", size=10, color=pal.muted, italic=True,
        )


def primitive_dropcap_paragraph(slide, paragraph: str, kicker: str | None = None, style: str = "editorial"):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    if kicker:
        add_text(
            slide, kicker.upper(),
            x_in=1.5, y_in=0.8, w_in=10.3, h_in=0.4,
            font="Inter", size=11, color=pal.accent, tracking=200,
        )

    first_char = paragraph[0]
    rest = paragraph[1:]

    add_text(
        slide, first_char,
        x_in=1.5, y_in=1.4, w_in=1.2, h_in=2.2,
        font=typ.headline_font, size=110, color=pal.accent, line_spacing=0.9,
    )
    add_text(
        slide, rest,
        x_in=2.7, y_in=1.5, w_in=8.5, h_in=5.5,
        font=typ.body_font, size=14, color=pal.ink, line_spacing=1.6,
    )


def primitive_clipped_headline(slide, text: str, style: str = "brutalist", clip_side: str = "right"):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    # Place oversized text starting at edge so the end is clipped
    width_in = 18 if clip_side == "right" else 15
    x = -1.0 if clip_side == "left" else 0.3
    add_text(
        slide, text,
        x_in=x, y_in=1.5, w_in=width_in, h_in=4.5,
        font=typ.headline_font, size=180, color=pal.ink, line_spacing=0.95,
    )


def primitive_text_as_image(slide, word: str, style: str = "brutalist"):
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)
    add_text(
        slide, word.upper(),
        x_in=0.0, y_in=1.0, w_in=14.0, h_in=6.0,
        font=typ.headline_font, size=320, color=pal.accent, line_spacing=0.9,
    )


def primitive_mono_callout(slide, lines: Sequence[str], position: str = "left", style: str = "brutalist"):
    tokens = STYLES[style]
    pal = tokens["palette"]
    set_background(slide, pal.bg)

    x = 0.8 if position == "left" else 8.0
    body = "\n".join(lines)
    add_text(
        slide, body,
        x_in=x, y_in=2.0, w_in=5.5, h_in=4.0,
        font="Consolas", size=12, color=pal.accent, line_spacing=1.5,
    )


def primitive_timeline(slide, nodes: Sequence[tuple[str, str, str]], style: str = "swiss"):
    """nodes: [(date, heading, blurb), ...]"""
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    n = len(nodes)
    if n == 0:
        return
    x_start, x_end = 0.8, 12.5
    y = 4.0
    add_line(slide, x_start, y, x_end, y, color=pal.ink, width_pt=1.2)
    step = (x_end - x_start) / max(n - 1, 1)
    for i, (date, heading, blurb) in enumerate(nodes):
        cx = x_start + i * step
        # Dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - 0.08), Inches(y - 0.08), Inches(0.16), Inches(0.16)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(pal.accent)
        dot.line.fill.background()
        # Date
        add_text(
            slide, date, x_in=cx - 0.8, y_in=y + 0.25, w_in=1.6, h_in=0.4,
            font=typ.body_font, size=11, color=pal.muted, align="center",
        )
        # Heading
        add_text(
            slide, heading, x_in=cx - 1.2, y_in=y + 0.7, w_in=2.4, h_in=0.5,
            font=typ.headline_font, size=14, color=pal.ink, align="center",
        )
        # Blurb
        add_text(
            slide, blurb, x_in=cx - 1.2, y_in=y + 1.2, w_in=2.4, h_in=1.5,
            font=typ.body_font, size=10, color=pal.ink, align="center", line_spacing=1.3,
        )


def primitive_half_bleed_placeholder(slide, caption: str, side: str = "right", style: str = "editorial"):
    """Reserved-image variant: paints a colored block where the image would go.
    Recipes can swap this for slide.shapes.add_picture()."""
    tokens = STYLES[style]
    pal, typ = tokens["palette"], tokens["type"]
    set_background(slide, pal.bg)

    if side == "right":
        add_rect(slide, 6.667, 0, 6.667, 7.5, fill=pal.muted)
        add_text(
            slide, caption,
            x_in=0.8, y_in=3.0, w_in=5.5, h_in=2.0,
            font=typ.headline_font, size=36, color=pal.ink, line_spacing=1.15,
        )
    else:
        add_rect(slide, 0, 0, 6.667, 7.5, fill=pal.muted)
        add_text(
            slide, caption,
            x_in=7.2, y_in=3.0, w_in=5.5, h_in=2.0,
            font=typ.headline_font, size=36, color=pal.ink, line_spacing=1.15,
        )


# ----- Presentation factory -----------------------------------------------

def new_presentation():
    """Create a 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # "Blank"
    return prs.slides.add_slide(blank_layout)
